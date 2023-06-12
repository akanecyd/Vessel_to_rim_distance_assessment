from pptx import Presentation
from pptx.util import Inches
from pptx.chart.data import ChartData, CategoryChartData, XyChartData, BubbleChartData

from pptx.util import Cm, Pt


def add_title(slide, msg):
    """
    args
        slide[slide]: Slide object
        msg[str] : Slide title message
    return:
        None
    """
    shapes = slide.shapes
    shapes.title.text = msg

def add_text(slide, msg, left, top, width, height, font_size, is_bold):
    """
    args:
        slide[slide]: Slide object
        msg[str]: Text box message
        left[int]: Position from the left end
        top[int] : Position from top
        width[int]: Width of object
        height[int]: Height of object
        font_size[int]: Font size
        is_bold[int]: Whether to make the text bold
    return:
        None
    """
    textbox = slide.shapes.add_textbox(Cm(left), Cm(top), Cm(width), Cm(height))
    p = textbox.text_frame.add_paragraph()
    p.text = msg
    p.font.size = Pt(font_size)
    p.font.bold = is_bold

def add_img(slide, img_path, left, top, width, height):
    """
    args:
        slide[slide]: Slide object
        img_path[str] : Image file path
        left[int]: Position from the left end
        top[int] : Position from top
        width[int]: Width of object
        height[int]: Height of object
    return:
        None
    """
    pic = slide.shapes.add_picture(img_path, 0, 0)
    pic.width = Cm(width)
    pic.height = Cm(height)
    pic.left = Cm(left)
    pic.top = Cm(top)